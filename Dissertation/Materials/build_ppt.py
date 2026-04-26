import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

NAVY   = RGBColor(0x0D, 0x2B, 0x45)
BLUE   = RGBColor(0x1E, 0x6F, 0xA8)
ACCENT = RGBColor(0x2A, 0xB4, 0xC4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF4, 0xF6, 0xF9)
DGREY  = RGBColor(0x44, 0x44, 0x55)
GREEN  = RGBColor(0x1A, 0x7A, 0x5E)
PURP   = RGBColor(0x7B, 0x2D, 0x8B)

blank = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    return sh

def txt(slide, text, l, t, w, h, size=18, bold=False,
        color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def bullets(slide, lines, l, t, w, h, size=15, color=DGREY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def top_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.25, NAVY)
    rect(slide, 0, 1.25, 13.33, 0.06, ACCENT)
    rect(slide, 0, 1.31, 13.33, 6.19, LGREY)
    txt(slide, title, 0.35, 0.12, 12.5, 0.9,
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.9, 12.5, 0.35,
            size=13, color=ACCENT)

def card(slide, l, t, w, h, head, head_col=BLUE):
    rect(slide, l, t, w, h, WHITE)
    rect(slide, l, t, w, 0.38, head_col)
    txt(slide, head, l+0.12, t+0.05, w-0.2, 0.3,
        size=13, bold=True, color=WHITE)

def two_col(slide, lhead, rhead, left_lines, right_lines,
            top=1.55, h=5.6, lhcol=BLUE, rhcol=BLUE):
    W = 6.1
    card(slide, 0.3, top, W, h, lhead, lhcol)
    bullets(slide, left_lines, 0.45, top+0.45, W-0.25, h-0.55)
    card(slide, 6.93, top, W, h, rhead, rhcol)
    bullets(slide, right_lines, 7.08, top+0.45, W-0.25, h-0.55)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1  Title
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 5.0, 7.5, NAVY)
rect(s, 5.0, 0, 0.08, 7.5, ACCENT)
rect(s, 5.08, 0, 8.25, 7.5, WHITE)

txt(s, "Psychological\nNarrative Engine", 0.35, 0.9, 4.4, 2.8,
    size=34, bold=True, color=WHITE)
txt(s, "PNE", 0.35, 3.7, 4.4, 0.9,
    size=52, bold=True, color=ACCENT)
txt(s, "Jerome Bawa  |  CS3IP Individual Project", 0.35, 5.4, 4.4, 0.45,
    size=12, color=LGREY)
txt(s, "University of Reading  |  2025-26", 0.35, 5.82, 4.4, 0.4,
    size=12, color=LGREY)

rect(s, 5.6, 1.2, 7.4, 5.5, WHITE)
rect(s, 5.6, 1.2, 7.4, 0.45, BLUE)
txt(s, "Project at a Glance", 5.75, 1.25, 7.1, 0.38,
    size=15, bold=True, color=WHITE)
bullets(s, [
    "Problem:  NPC dialogue in games is scripted, static, psychologically hollow",
    "",
    "Solution:  A BDI pipeline that makes NPCs think before they speak",
    "",
    "How:  Cognition is deterministic (no AI) -> local AI speaks the result",
    "",
    "What:  Beliefs, Desires, Intentions -- no pre-written lines needed",
    "",
    "Where:  Game-engine agnostic REST + WebSocket API",
    "         Plug into Unity, Godot, Unreal Engine",
    "",
    "Who:  Indie and mid-scale developers on consumer hardware (4-6 GB VRAM)",
], 5.75, 1.72, 7.1, 4.85, size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2  Introduction
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Introduction  -  The Problem",
        "Chapter 1: Why dialogue trees are not enough")

rect(s, 0.3, 1.48, 12.73, 1.12, WHITE)
rect(s, 0.3, 1.48, 0.09, 1.12, ACCENT)
txt(s,
    '"NPCs in scripted systems do not hold beliefs or form intentions -- they evaluate '
    'conditions and return pre-written strings. The apparent depth of a well-written '
    'branching tree is purely a product of human authorship, not of any computational '
    'model of mind."',
    0.52, 1.53, 12.3, 1.0, size=13, italic=True, color=DGREY)

two_col(s,
    "The Status Quo", "The Research Question",
    [
        "- The dialogue tree has dominated NPC interaction for 30+ years",
        "",
        "- Every response is pre-authored -- developers write every word",
        "",
        "- Scaling is expensive: more NPCs = proportionally more writing,",
        "  QA, and voice acting",
        "",
        "- NPCs pattern-match player choices -- they do not reason",
        "",
        "- Disco Elysium (2019) showed emergent narrative is possible",
        "  -- but only on the player side",
    ],
    [
        "- Can an NPC hold an internal model of the world and respond from it?",
        "",
        "- Can psychological depth and player agency coexist?",
        "",
        "- Can cognition be separated from text generation?",
        "",
        "The PNE answers: Yes",
        "",
        "-> Build a cognitive pipeline that thinks",
        "   then let a local AI speak the result",
        "",
        "Target: indie/mid-scale developers, consumer hardware",
    ],
    top=2.7, h=4.55)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3  Literature Review
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Literature Review",
        "Chapter 3: Four bodies of work that shape the PNE")

CW = 3.05; GAP = 0.13; TOP = 1.52; CH = 5.65
cols = [0.3, 0.3+CW+GAP, 0.3+(CW+GAP)*2, 0.3+(CW+GAP)*3]

for l, hcol, head, lines in [
    (cols[0], NAVY, "BDI Agent Architecture", [
        "Bratman (1987)",
        "Beliefs -> Desires -> Intentions",
        "Foundation of the entire pipeline",
        "",
        "Rao & Georgeff (1995)",
        "BDI from theory to practice",
        "",
        "Georgeff et al. (1999)",
        "Formal BDI model of agency",
        "",
        "-> Every pipeline stage maps directly",
        "   to BDI theory",
    ]),
    (cols[1], BLUE, "Cognitive Psychology", [
        "Ellis (1962)",
        "A-B-C: events filtered through beliefs",
        "-> Cognitive bias library",
        "",
        "Neisser (1967)",
        "Cognition is constructive, not retrieval",
        "-> Thought-pattern templates",
        "",
        "Heider (1958)",
        "Attribution theory + locus of control",
        "-> NPC personality parameter",
    ]),
    (cols[2], GREEN, "Social Psychology", [
        "Cialdini (1984)",
        "Six principles of persuasion",
        "-> Tone scoring system",
        "",
        "Davis (1983)",
        "Empathy as measurable dimension",
        "-> Empathy NPC attribute",
        "",
        "Festinger (1954)",
        "Social comparison & consistency",
        "-> Desire formation grounding",
    ]),
    (cols[3], PURP, "Narrative & Prior Systems", [
        "Mateas & Stern (2003)",
        "Facade -- interactive drama",
        "Closest prior system; compared directly",
        "",
        "Park et al. (2023)",
        "Generative Agents (LLM NPCs)",
        "Modern benchmark",
        "",
        "Murray (1997)",
        "Hamlet on the Holodeck",
        "",
        "Gap: none combine structured",
        "psychology with a local LLM",
    ]),
]:
    card(s, l, TOP, CW, CH, head, hcol)
    bullets(s, lines, l+0.12, TOP+0.45, CW-0.22, CH-0.55, size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4  Methodology
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Methodology",
        "Chapter 4: How the system was designed and tested")

two_col(s,
    "Development Approach", "Evaluation Strategy",
    [
        "- Agile iterative development (logbook-documented sprints)",
        "",
        "- Each component built, tested, then integrated",
        "",
        "- 3 representative NPCs -- contrasting personality profiles:",
        "   Krakk  (high empathy, low assertion)",
        "   Moses  (max assertion, near-zero cognitive flexibility)",
        "   Troy   (low self-esteem, high assertion, zero independence)",
        "",
        "- Same scenario across all three: door_guard_night",
        "",
        "- No user studies -- internal state inspection as primary validation",
        "",
        "- Ethical: psychology constructs are design parameters,",
        "  not clinical representations",
    ],
    [
        "Three evaluation questions:",
        "",
        "1. Does internal NPC state stay consistent turn to turn?",
        "",
        "2. Does the AI output match what the pipeline computed?",
        "",
        "3. Does the FSM route to the correct narrative outcome?",
        "",
        "- Cognitive layer tested independently (no AI needed)",
        "  -- deterministic, fully unit-testable",
        "",
        "- Full pipeline trace logged per conversation turn",
        "",
        "- Comparative player build test (Troy: Empathy vs Assertion)",
        "  to verify player agency is real",
    ],
    top=1.52, h=5.65)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5  Architecture -- Pipeline image
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Architecture  -  The PNE Pipeline",
        "Chapter 5: Seven stages from player input to NPC speech")

img_path = r"d:/Programming/PNE (Github)/cs3ip/Dissertation/Materials/PNE Pipeline.png"
try:
    s.shapes.add_picture(img_path, Inches(0.3), Inches(1.42), Inches(8.9), Inches(5.85))
except Exception as e:
    txt(s, "[Pipeline image could not load: " + str(e) + "]",
        0.3, 1.5, 8.9, 5.5, size=12, color=DGREY)

card(s, 9.35, 1.42, 3.68, 5.85, "Stage Summary", NAVY)
bullets(s, [
    "I    NPC Intent Layer",
    "     Loads identity, goals, terminal endings",
    "",
    "II   Choice Selection",
    "     Player picks; tone is scored",
    "",
    "IIa  Skill Check",
    "     2 biased dice: player vs NPC",
    "",
    "III  Cognitive Interpretation",
    "     NPC forms an internal belief",
    "",
    "IV   Desire Formation",
    "     Belief -> goal type",
    "",
    "V    Socialisation Filter",
    "     Goal -> named behaviour",
    "",
    "VI   Conversational Output",
    "     Ollama generates the spoken text",
], 9.5, 1.87, 3.42, 5.28, size=12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6  Architecture -- Key Innovations
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Architecture  -  Key Innovations",
        "Chapter 5: What makes the PNE different")

CW2 = 6.15; CH2 = 2.5; GAP2 = 0.23
positions = [
    (0.3,        1.5),
    (0.3+CW2+GAP2, 1.5),
    (0.3,        1.5+CH2+GAP2),
    (0.3+CW2+GAP2, 1.5+CH2+GAP2),
]
innovations = [
    ("CognitiveThoughtMatcher  --  no AI in the reasoning layer", NAVY, [
        "810 scored thought templates matched against player tone + NPC personality",
        "Produces: internal thought, belief, emotional reaction, cognitive bias",
        "Fully deterministic and unit-testable -- AI kept entirely out of this layer",
        "Two NPCs with different profiles -> provably different thoughts, same input",
    ]),
    ("2-Dice Probabilistic Skill Check", BLUE, [
        "Player + NPC each roll a biased d6 -- player wins on a tie",
        "Bias derived from player skill level and NPC resistance threshold",
        "success_pct shown before each choice so players can weigh risk",
        "Failed paths permanently pruned -- raises genuine, lasting stakes",
    ]),
    ("Judgement-Score Finite State Machine", GREEN, [
        "0-100 scalar aggregated across all turns -- routes narrative independently of AI",
        "Risk multiplier: bold low-probability choices carry amplified consequences",
        "FSM uses judgement score + relation score -- not raw AI text classification",
        "Narrative destination is predictable and author-controlled",
    ]),
    ("Game-Engine Agnostic API", PURP, [
        "FastAPI REST + WebSocket -- Unity, Godot, Unreal can all connect",
        "AI tokens streamed word-by-word via WebSocket for real-time feel",
        "Drop-in C# client (Unity) and GDScript client (Godot) provided",
        "No cloud dependency -- fully offline, zero per-request cost",
    ]),
]
for (l, t), (head, hcol, lines) in zip(positions, innovations):
    card(s, l, t, CW2, CH2, head, hcol)
    rect(s, l, t, 0.08, CH2, ACCENT)
    bullets(s, lines, l+0.2, t+0.45, CW2-0.3, CH2-0.55, size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7  Results
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Results",
        "Chapter 6: Three NPCs, same scenario  --  door_guard_night")

# Table
rows = [
    ("NPC",        "Personality Profile",                              "Opening",    "Turns", "Result"),
    ("Krakk",      "High empathy (0.8), low assertion (0.3)\nInferiority wildcard", "Diplomatic", "2", "SUCCEED"),
    ("Moses",      "Max assertion (1.0), low cog. flex (0.3)\nMartyr wildcard x2",  "Authority",  "4", "SUCCEED"),
    ("Troy",       "Low self-esteem (0.2), high assertion (0.8)\nNo wildcard",       "Authority",  "4", "FAIL"),
    ("Troy (alt)", "Same NPC -- Empathy build player skills",          "Empathy",    "2", "SUCCEED"),
]
col_ws = [1.45, 4.55, 1.8, 0.85, 1.45]
col_ls = [0.3]
for cw in col_ws[:-1]:
    col_ls.append(col_ls[-1]+cw)
RH = 0.68
for ri, row in enumerate(rows):
    bg = NAVY if ri == 0 else (WHITE if ri % 2 == 1 else LGREY)
    total_w = sum(col_ws)
    rect(s, 0.3, 1.48 + ri*RH, total_w, RH, bg)
    for ci, (cell, cw, cl) in enumerate(zip(row, col_ws, col_ls)):
        c = WHITE if ri == 0 else DGREY
        if ri > 0 and ci == 4:
            c = GREEN if "SUCCEED" in cell else RGBColor(0xCC, 0x22, 0x22)
        txt(s, cell, cl+0.08, 1.48+ri*RH+0.08, cw-0.12, RH-0.08,
            size=12, bold=(ri == 0), color=c)

# Three callout boxes
bx = [(0.3, "Pipeline Coherence"), (4.55, "Player Agency"), (8.8, "Coherence Gap")]
bw = 4.0
for bx_l, bx_head in bx:
    rect(s, bx_l, 5.0, bw, 2.22, WHITE)
    rect(s, bx_l, 5.0, bw, 0.32, ACCENT)
    txt(s, bx_head, bx_l+0.15, 5.02, bw-0.2, 0.28, size=12, bold=True, color=WHITE)

bullets(s, [
    "Internal state consistent across all turns for all 3 NPCs",
    "Cognitive bias correctly shaped desire type every case",
    "Moses Martyr wildcard fired and held across turns",
], 0.45, 5.35, bw-0.25, 1.75, size=12)

bullets(s, [
    "Troy Empathy build -> SUCCEED in 2 turns",
    "Troy Assertion build -> FAIL at turn 3",
    "Outcome driven by skill distribution, not choice text",
], 4.7, 5.35, bw-0.25, 1.75, size=12)

bullets(s, [
    "Moses turn 3: pipeline said explosive / confrontation 0.978",
    "LLM produced a more moderate response",
    "Pipeline logic correct -- LLM did not express full intensity",
], 8.95, 5.35, bw-0.25, 1.75, size=12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8  Discussion & Analysis
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Discussion & Analysis",
        "Chapter 7: What the results tell us -- and where the limits are")

two_col(s,
    "What Worked", "Limitations & Honest Critique",
    [
        "BDI pipeline produces psychologically coherent NPC behaviour",
        "without requiring the AI to reason",
        "",
        "Separating cognition from language generation works:",
        "  reasoning is deterministic and testable;",
        "  speech is constrained by 19 canonical intentions",
        "",
        "Personality parameters meaningfully differentiate NPCs --",
        "Krakk, Moses and Troy behave exactly as their profiles predict",
        "",
        "Player agency is real: skill distribution determines terminal",
        "routing, not just which words the player chose",
        "",
        "Local AI (Qwen2.5:3b) is viable for real-time dialogue on",
        "consumer hardware",
        "",
        "Research question answered: Yes -- psychological realism,",
        "narrative coherence, and player agency can coexist",
    ],
    [
        "AI coherence gap: model under-expressed intensity",
        "on Moses turn 3 despite identical pipeline state",
        "",
        "Wildcard system partially implemented -- complex",
        "multi-turn wildcard interactions not fully tested",
        "",
        "Logging gap: some internal transitions not persisted,",
        "limiting post-hoc analysis depth",
        "",
        "Narrow scenario coverage -- one scenario, three NPCs;",
        "generalisation claims are bounded",
        "",
        "4-6 GB VRAM requirement excludes low-spec hardware",
        "and mobile devices entirely",
        "",
        "Future validation needed across broader scenario types",
        "and larger NPC casts",
    ],
    top=1.52, h=5.65,
    lhcol=GREEN, rhcol=RGBColor(0xB0, 0x40, 0x40))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9  Conclusion & Future Work
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "Conclusion & Future Work",
        "Chapter 8: What was achieved -- and where next")

two_col(s,
    "Conclusion", "Future Work",
    [
        "The PNE demonstrates that a BDI pipeline + psychology-grounded",
        "personality model + constrained local AI is a viable architecture",
        "for adaptive NPC dialogue",
        "",
        "NPCs can think deterministically and speak naturally --",
        "without pre-authored lines and without cloud infrastructure",
        "",
        "The system is game-engine agnostic, open-source, and",
        "deployable on mid-range consumer hardware",
        "",
        "Key contribution: the CognitiveThoughtMatcher",
        "An AI-free, testable cognitive interpretation layer",
        "grounded in Ellis (1962), Heider (1958), Davis (1983)",
        "",
        "Practical impact: enables indie developers to author",
        "psychologically rich NPC dialogue at a fraction of",
        "traditional authorial cost",
    ],
    [
        "Multi-level mentalising",
        "  NPCs that model what the player believes about the NPC",
        "",
        "Persistent NPC memory",
        "  Relationship history that carries across sessions",
        "",
        "Multi-NPC interaction",
        "  NPCs that react to each other, not just to the player",
        "",
        "Expanded scenario coverage",
        "  Full story arcs, not just single-node scenarios",
        "",
        "Larger AI evaluation",
        "  Quantify coherence gap across model sizes (3B -> 13B)",
        "",
        "Player-facing transparency",
        "  UI that shows NPC emotional state in real time",
    ],
    top=1.52, h=5.65,
    lhcol=NAVY, rhcol=BLUE)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10  References
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
top_bar(s, "References", "Key sources cited in this presentation")

bullets(s, [
    "Bratman, M. (1987) Intention, plans, and practical reason.",
    "  Cambridge, MA: Harvard University Press.",
    "",
    "Rao, A.S. & Georgeff, M.P. (1995) 'BDI agents: from theory to practice'.",
    "  Proc. ICMAS-95. Menlo Park: AAAI Press, pp. 312-319.",
    "",
    "Georgeff, M. et al. (1999) 'The BDI model of agency'.",
    "  Intelligent Agents V, LNCS 1555. Berlin: Springer, pp. 1-10.",
    "",
    "Mateas, M. & Stern, A. (2003) 'Facade'.",
    "  Proc. Game Developers Conference, San Jose, CA.",
    "",
    "Park, J.S. et al. (2023) 'Generative agents'.",
    "  Proc. UIST '23. New York: ACM.",
], 0.3, 1.48, 6.5, 5.85, size=13)

bullets(s, [
    "Ellis, A. (1962) Reason and emotion in psychotherapy.",
    "  New York: Lyle Stuart.",
    "",
    "Heider, F. (1958) The psychology of interpersonal relations.",
    "  New York: Wiley.",
    "",
    "Davis, M.H. (1983) 'Measuring individual differences in empathy'.",
    "  J. Personality & Social Psychology, 44(1), pp. 113-126.",
    "",
    "Cialdini, R.B. (1984) Influence: the psychology of persuasion.",
    "  New York: Harper Collins.",
    "",
    "Murray, J.H. (1997) Hamlet on the Holodeck.",
    "  New York: Free Press.",
], 6.95, 1.48, 6.1, 5.85, size=13)

# ── Save ─────────────────────────────────────────────────────────────
out = r"d:/Programming/PNE (Github)/cs3ip/Dissertation/Materials/PNE_Presentation.pptx"
prs.save(out)
print("Saved:", out)
