from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours ──────────────────────────────────────────────────────────────────
C_BG         = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
C_SECTION_BG = RGBColor(0x1A, 0x2B, 0x4A)   # dark navy
C_ACCENT     = RGBColor(0xF5, 0xA6, 0x23)   # amber
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT      = RGBColor(0xD0, 0xD8, 0xEE)   # pale blue-white
C_BODY       = RGBColor(0xE8, 0xEC, 0xF5)   # near-white

CATEGORY_COLOURS = {
    "BDI Models":           RGBColor(0x2E, 0x5F, 0xC1),
    "Cognitive Psychology": RGBColor(0x1A, 0x80, 0x7A),
    "Social Psychology":    RGBColor(0x6B, 0x3A, 0xB0),
    "Narrative Design":     RGBColor(0xC1, 0x4B, 0x2E),
    "Local LLMs":           RGBColor(0x2E, 0x8B, 0x4A),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, colour):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, colour=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"
    return txb

def add_bullet_slide(prs, category, paper_title, author_year, slide_title,
                     bullets, label, accent_col):
    """Generic 3-column-free bullet slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)

    # Left accent bar
    add_rect(slide, 0, 0, 0.12, 7.5, accent_col)

    # Top strip
    add_rect(slide, 0.12, 0, 13.21, 1.45, C_SECTION_BG)

    # Category + label tag
    add_text(slide, f"{category}  ·  {label}", 0.28, 0.08, 9, 0.4,
             9, colour=C_ACCENT, bold=True)

    # Paper title in header
    add_text(slide, paper_title, 0.28, 0.42, 10.5, 0.55, 14,
             bold=True, colour=C_WHITE)

    # Author/year
    add_text(slide, author_year, 0.28, 0.92, 8, 0.38, 10,
             colour=C_LIGHT, bold=False)

    # Slide title
    add_text(slide, slide_title, 0.28, 1.55, 12.7, 0.55, 20,
             bold=True, colour=accent_col)

    # Divider line
    add_rect(slide, 0.28, 2.18, 12.55, 0.03, accent_col)

    # Bullets
    y = 2.32
    for bullet in bullets:
        # Bullet dot
        add_rect(slide, 0.28, y + 0.11, 0.08, 0.08, accent_col)
        add_text(slide, bullet, 0.48, y, 12.35, 0.68, 13,
                 colour=C_BODY, wrap=True)
        y += 0.82

    return slide


def add_section_slide(prs, category, accent_col, description):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
    add_rect(slide, 0, 0, 13.33, 0.08, accent_col)
    add_rect(slide, 0, 7.42, 13.33, 0.08, accent_col)
    add_rect(slide, 0, 0, 0.35, 7.5, accent_col)

    add_text(slide, "SECTION", 0.65, 2.2, 12, 0.5, 14,
             bold=True, colour=accent_col, align=PP_ALIGN.LEFT)
    add_text(slide, category, 0.65, 2.75, 12, 1.1, 44,
             bold=True, colour=C_WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, 0.65, 3.95, 4.5, 0.06, accent_col)
    add_text(slide, description, 0.65, 4.15, 11.8, 1.5, 14,
             colour=C_LIGHT, align=PP_ALIGN.LEFT, wrap=True)
    return slide


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)
    add_rect(slide, 0, 0, 13.33, 0.1, C_ACCENT)
    add_rect(slide, 0, 7.4, 13.33, 0.1, C_ACCENT)

    add_text(slide, "PSYCHOLOGICAL NARRATIVE ENGINE", 0.8, 1.8, 12, 0.7,
             18, bold=True, colour=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, "Literature Review Summary", 0.8, 2.55, 12, 1.0,
             38, bold=True, colour=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.5, 3.75, 6.33, 0.06, C_ACCENT)
    add_text(slide,
             "A plain-English guide to each source:\nwhat it says, how it connects to PNE, and what to write in Chapter 2.",
             0.8, 3.95, 12, 1.2, 14,
             colour=C_LIGHT, align=PP_ALIGN.CENTER, wrap=True)
    add_text(slide, "Jerome Bawa  ·  CS3IP Dissertation  ·  2026",
             0.8, 6.6, 12, 0.5, 11,
             colour=C_LIGHT, align=PP_ALIGN.CENTER)


# ── Paper data ────────────────────────────────────────────────────────────────
PAPERS = [
    # ── BDI Models ──────────────────────────────────────────────────────────
    {
        "category": "BDI Models",
        "title": "Reasoning about Rational Agents",
        "author_year": "Michael Wooldridge  ·  2000  ·  MIT Press",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "This book is the go-to textbook for understanding how intelligent agents (software programs that act on their own) make decisions.",
                    "It explains the BDI model: agents hold Beliefs (what they know), Desires (what they want), and Intentions (what they commit to doing).",
                    "Think of it like modelling a person's thought process: you believe it's raining, you desire to stay dry, so you intend to grab an umbrella.",
                    "Wooldridge shows how to turn this philosophy into working code — agents that reason step-by-step about their goals and update their plans when the world changes.",
                    "It also compares BDI to other agent types (reactive agents that just respond, logic-based agents that reason from rules) and explains trade-offs.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "PNE's entire processing pipeline is a BDI loop: each player turn triggers belief update → desire formation → intention selection.",
                    "Wooldridge's formalism justifies WHY this three-stage structure is the right model for NPCs that need to 'think' rather than just pattern-match.",
                    "The book's discussion of intention persistence (agents don't drop goals at the first obstacle) maps to PNE's recovery mechanic — NPCs don't flip immediately on one failed roll.",
                    "It provides the academic authority to say 'PNE is a BDI agent system' rather than just 'PNE uses if-then rules with extra steps'.",
                    "Also useful for contrasting: Wooldridge's agents use formal logic; PNE uses template-matching — a deliberate pragmatic trade-off this book lets you articulate.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Open by citing Wooldridge as the foundational definition of BDI: beliefs, desires, and intentions as distinct representational states that drive rational action.",
                    "Use it to explain WHY BDI is the right cognitive model for NPCs: it captures goal-directedness, context-sensitivity, and the ability to revise plans.",
                    "Contrast: Wooldridge's theoretical BDI agents use formal logic proofs; PNE uses weighted template-matching. Acknowledge this and argue it's a practical design choice for real-time game dialogue.",
                    "Reference in the paragraph that introduces BDI before discussing Rao & Georgeff and Bratman — Wooldridge synthesises both into an accessible framework.",
                ],
            },
        ],
    },
    {
        "category": "BDI Models",
        "title": "Programming Multi-Agent Systems in AgentSpeak using Jason",
        "author_year": "Bordini, Hübner & Wooldridge  ·  2007  ·  Wiley",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "This is the practical companion to BDI theory: a book showing you how to actually BUILD multi-agent systems using a language called AgentSpeak (implemented in the Jason framework).",
                    "AgentSpeak lets you write agents as a list of 'plans': if I believe X and want Y, then do Z. Agents run many plans in parallel and pick the best one for the situation.",
                    "Jason is used in research and teaching as a reference implementation of BDI — it's the closest thing to a 'canonical' BDI programming environment.",
                    "The book covers belief bases, goal-directed behaviour, plan libraries, and how multiple agents communicate and coordinate with each other.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "Jason stores agent behaviour in external plan libraries (files of if-then rules). PNE stores NPC cognition in cognitive_thoughts.json and the intention registry — a similar concept, different implementation.",
                    "This is the key CONTRAST: Jason agents select plans dynamically at runtime from a formal library. PNE selects templates via weighted matching — simpler, faster, easier for game authors to edit.",
                    "Using this as a contrast lets you argue that PNE makes a deliberate trade-off: less formal correctness, more authorial control and runtime speed — appropriate for real-time game dialogue.",
                    "Also useful for showing PNE is aware of the field: you know what a 'proper' BDI implementation looks like, and you made conscious choices to deviate.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Jason as the primary contrast case for PNE's BDI implementation: 'Where established BDI frameworks such as Jason externalise agent reasoning in formal plan libraries, PNE internalises cognition through template-matching...'",
                    "This lets you show you understand the state of the art and have made deliberate, reasoned design choices rather than simply not knowing the alternatives.",
                    "Mention it after introducing Wooldridge's theory — Jason is theory in practice, which then positions PNE as a further pragmatic adaptation of that practice for game contexts.",
                    "Keep it brief: 3-5 sentences. Its main job is contrast, not detailed description.",
                ],
            },
        ],
    },
    {
        "category": "BDI Models",
        "title": "The Belief-Desire-Intention Model of Agency",
        "author_year": "Georgeff, Pell, Pollack, Tambe & Wooldridge  ·  1999  ·  Intelligent Agents V, LNCS 1555",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Paper About?",
                "bullets": [
                    "A short, dense retrospective written by the people who actually built the BDI model. It reviews 10+ years of BDI research from theory to working implementations.",
                    "They explain the gap between Bratman's philosophical BDI (1987 — a theory of human planning) and the computational BDI used in AI (Rao & Georgeff, 1995).",
                    "The paper argues that BDI isn't just one thing: there are many variants, and different applications need different versions. Some emphasise belief revision, others commitment, others communication.",
                    "It reads as a 'state of the art' snapshot: here's what BDI can do, here's what's still hard, here's where it should go.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "This paper is the bridge between Bratman (philosophy) and Rao & Georgeff (formalisation) — both already cited in Chapter 1. It contextualises both in one place.",
                    "The authors' point that BDI has many variants is directly relevant: PNE is another variant, adapted for NPC dialogue rather than autonomous robots or logic-based planning.",
                    "Their discussion of intention commitment maps to PNE's design: NPCs maintain an intention through a turn (they don't abandon mid-sentence), but update between turns.",
                    "Validates the choice to use BDI as a framework at all: if the architects of BDI say it generalises across domains, that supports applying it to NPC conversation.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use this as the citation when you describe the evolution from Bratman's philosophy to computational BDI: 'Georgeff et al. (1999) trace this development, noting that BDI has been instantiated in multiple computational frameworks...'",
                    "It's particularly useful for the sentence that acknowledges BDI variants exist and positions PNE within that landscape.",
                    "Works well as a secondary citation alongside Rao & Georgeff (1995) — they're companions. Don't need more than a paragraph; its role is intellectual lineage.",
                    "Good for the opening of your BDI section in Chapter 2 before you get into Wooldridge's framework in detail.",
                ],
            },
        ],
    },

    # ── Cognitive Psychology ─────────────────────────────────────────────────
    {
        "category": "Cognitive Psychology",
        "title": "Cognitive Psychology",
        "author_year": "Ulric Neisser  ·  1967  ·  Appleton-Century-Crofts",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "The book that invented the field. Neisser coined the term 'cognitive psychology' and established that the mind actively constructs its understanding of the world — it doesn't just passively receive information.",
                    "Central idea: SCHEMA THEORY. We don't experience reality directly; we interpret it through mental frameworks (schemas) built from past experience. Two people can see the same event very differently.",
                    "He introduces the idea that perception, memory, and thinking are all constructive processes — the brain fills in gaps, makes assumptions, and applies patterns.",
                    "This was radical in 1967 because psychology before it was behaviourist: only observable behaviour mattered. Neisser opened the 'black box' of the mind.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "Schema theory is the direct theoretical basis for PNE's cognitive interpretation layer: NPCs don't hear the player's words neutrally — they interpret them through their cognitive schema (self_esteem, locus_of_control, cog_flexibility).",
                    "The 810 templates in cognitive_thoughts.json are literally schema-driven interpretations: the same player input produces different NPC thoughts depending on which schema is active.",
                    "Neisser's constructivism justifies WHY NPCs with different cognitive profiles should generate different beliefs from identical input — it mirrors how humans actually work.",
                    "Gives academic grounding to the design decision to make cognition template-based rather than neutral or random.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Neisser to introduce the cognitive psychology section: 'Neisser (1967) established that human perception is an active, constructive process mediated by cognitive schemas — pre-existing frameworks that shape interpretation of new information.'",
                    "Then connect to PNE: 'This constructivist model provides the theoretical basis for PNE's cognitive interpretation layer, in which NPC personality parameters modulate the template selected to interpret player input.'",
                    "Keep it foundational — 2-3 sentences. Its job is to give the cognitive section its theoretical footing before you get into Beck and Ellis.",
                    "Positions PNE's design as grounded in decades of cognitive science, not invented from scratch.",
                ],
            },
        ],
    },
    {
        "category": "Cognitive Psychology",
        "title": "Reason and Emotion in Psychotherapy",
        "author_year": "Albert Ellis  ·  1962  ·  Lyle Stuart",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "Ellis created Rational Emotive Behaviour Therapy (REBT) and this is its founding text. Core idea: it's not what happens to you that causes distress — it's what you BELIEVE about what happens.",
                    "He introduces the A-B-C model: A (Activating event) → B (Belief about it) → C (emotional Consequence). Same event, different belief, completely different emotional outcome.",
                    "Ellis catalogued irrational beliefs — systematic errors in thinking that lead people to feel and behave in self-defeating ways. Things like 'if someone disagrees with me, they must hate me.'",
                    "The book is practical and clinical, but the underlying model is deeply cognitive: what matters is the interpretive step between stimulus and response.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "The A-B-C model IS the PNE cognitive pipeline: player input (A) → NPC belief via template-matching (B) → emotional valence + desire formation (C).",
                    "Ellis's irrational beliefs map directly to PNE's cognitive bias types: hostile attribution ('they must mean harm') is an irrational belief; cynical realism and projection are others.",
                    "The idea that the same event triggers different beliefs in different people validates having different cognitive templates for NPCs with different bias profiles.",
                    "Desire formation in PNE (information-seeking, affiliation, protection, dominance) mirrors Ellis's emotional consequences — what the NPC 'wants' emerges from what they 'believe'.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Cite Ellis when explaining the cognitive interpretation → desire pipeline: 'Ellis's A-B-C model (1962) provides a structural parallel: the player's dialogue choice constitutes the activating event (A); the NPC's template-matched thought constitutes the belief (B); the resulting desire state constitutes the emotional consequence (C).'",
                    "Use this to show the pipeline has a grounded psychological model behind it, not an arbitrary three-step sequence.",
                    "Also use Ellis when defining the cognitive bias types in cognitive_thoughts.json — his irrational belief categories are the intellectual ancestors of those bias labels.",
                    "3-5 sentences is enough. Its main role is structural analogy between REBT and the PNE pipeline.",
                ],
            },
        ],
    },
    {
        "category": "Cognitive Psychology",
        "title": "The Psychology of Interpersonal Relations",
        "author_year": "Fritz Heider  ·  1958  ·  Wiley",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "The book that founded attribution theory — the study of how people explain other people's behaviour. When someone does something, do you assume it's because of WHO THEY ARE or because of the SITUATION they're in?",
                    "Heider called these internal attribution (personality, intent) vs external attribution (circumstances, accident). Crucially: people tend to over-attribute to personality and under-attribute to situation (the 'fundamental attribution error').",
                    "He also introduced the concept of BALANCE: people feel uncomfortable when their beliefs about relationships are inconsistent (e.g. liking someone who dislikes your friend creates psychological tension).",
                    "The whole book is about how ordinary people make sense of social reality — their naive 'folk psychology' of intentions, reasons, and causes.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "Attribution theory is the engine behind locus_of_control in PNE: an NPC with high external locus attributes player behaviour to circumstances; high internal locus attributes it to the player's character and intent.",
                    "Hostile attribution bias (the most used bias type in cognitive_thoughts.json) is a direct application of Heider: the NPC assumes the player's intent is hostile regardless of the actual words used.",
                    "The balance theory maps to player_relation: relationship tension (low relation + positive player input) produces cognitive dissonance in the NPC — which the desire layer resolves by seeking 'information' or 'protection'.",
                    "Heider grounds the whole cognitive interpretation layer in decades of social cognition research.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Heider to introduce attribution as the mechanism behind locus_of_control: 'Heider (1958) established that individuals differ systematically in whether they attribute others' behaviour to internal disposition or external circumstance — a distinction formalised in PNE's locus_of_control parameter.'",
                    "Also cite when defining hostile attribution bias: it is Heider's internal attribution applied in its most extreme and distorted form.",
                    "Works best as a 3-4 sentence supporting citation in the cognitive psychology section, after Neisser and Ellis have established the schema/belief framework.",
                    "Gives PNE's personality parameters a specific theoretical lineage, not just a list of made-up sliders.",
                ],
            },
        ],
    },

    # ── Social Psychology ────────────────────────────────────────────────────
    {
        "category": "Social Psychology",
        "title": "Influence: The Psychology of Persuasion",
        "author_year": "Robert B. Cialdini  ·  1984  ·  Harper Collins",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "One of the most widely-read books in social psychology. Cialdini studied compliance professionals (salespeople, negotiators, advertisers) to find universal principles of persuasion.",
                    "He identified six principles: Reciprocity (you gave me something, so I owe you), Commitment (once I agree, I keep agreeing), Social Proof (others are doing it), Authority (experts say so), Liking (I'm more persuaded by people I like), and Scarcity (less available = more desirable).",
                    "The core insight: humans use mental shortcuts (heuristics) when deciding whether to comply with a request. These shortcuts can be triggered deliberately by framing the same request in different ways.",
                    "It's both a description of how persuasion works psychologically AND a practical guide to the techniques — making it unusually applicable to dialogue design.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "PNE's four Language Arts (authority, diplomacy, empathy, manipulation) map almost directly onto Cialdini's six principles: authority→Authority, diplomacy→Liking+Reciprocity, empathy→Liking, manipulation→Social Proof+Scarcity+Commitment.",
                    "This gives the Language Arts system a rigorous theoretical grounding: they're not arbitrary dialogue flavours, they're implementations of empirically validated persuasion mechanisms.",
                    "Cialdini's insight that the SAME request can succeed or fail depending on framing mirrors PNE's design: the same choice text at different tone_weights produces different NPC responses.",
                    "Also relevant to the social intention layer: NPC intentions like 'Assert Dominance' or 'Seek Connection' are Cialdini principles enacted from the NPC's side.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Cialdini to anchor the Language Arts system: 'PNE's four Language Arts — authority, diplomacy, empathy, and manipulation — are grounded in Cialdini's (1984) six principles of persuasion, which establish that compliance is reliably modulated by the rhetorical framing of a request rather than its content alone.'",
                    "This is a strong paragraph because it both cites a well-known work AND shows that PNE's design decisions weren't arbitrary.",
                    "Also reference when discussing NPC intention selection — the social layer chooses intentions that mirror Cialdini's compliance triggers from the NPC's perspective.",
                    "Around 4-6 sentences total across the social psychology section.",
                ],
            },
        ],
    },
    {
        "category": "Social Psychology",
        "title": "Social Comparison and Social Psychology",
        "author_year": "Serge Guimond (Ed.)  ·  2006  ·  Cambridge University Press",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "An edited academic volume collecting research on social comparison theory — the idea (originally Festinger, 1954) that humans evaluate their own opinions and abilities by comparing themselves to others.",
                    "The book extends this into cognition, intergroup relations, and cross-cultural contexts: not just 'am I better or worse than this person' but 'how does my GROUP compare to other groups, and how does that shape my identity?'",
                    "Key chapters cover how social comparisons drive self-esteem, how group membership shapes what comparisons feel relevant, and how culture moderates these effects.",
                    "Guimond also synthesises the link between social comparison and Social Identity Theory (Tajfel & Turner) — in-group favouritism emerges partly from self-evaluative comparisons.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "NPC self_esteem in PNE is partly a social comparison construct: how does the NPC rate themselves relative to the player, their faction, their rivals? Guimond provides the theoretical basis for this.",
                    "Faction dynamics in PNE (in-group/out-group tension when a player from outside the Insurgency talks to Troy) are grounded in the intergroup comparison mechanisms Guimond covers.",
                    "conf_indep (confidence/independence) is influenced by social comparison: an NPC who constantly compares themselves to stronger figures will have lower conf_indep, affecting intention selection.",
                    "Also supports the idea that player_relation isn't just a number — it represents the NPC's ongoing social self-evaluation relative to the player.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Guimond to ground self_esteem and faction-based NPC behaviour: 'Guimond (2006) establishes that self-evaluation is inherently comparative and socially situated — an individual's self-esteem is modulated by perceived standing relative to in-group and out-group members, a dynamic encoded in PNE's self_esteem and faction parameters.'",
                    "Keeps the social psychology section from being only about persuasion (Cialdini) — adds the dimension of identity and group membership.",
                    "3-4 sentences is sufficient. Position it after Cialdini, before Goffman — it bridges persuasion mechanisms and the performative self-presentation Goffman covers.",
                ],
            },
        ],
    },
    {
        "category": "Social Psychology",
        "title": "The Presentation of Self in Everyday Life",
        "author_year": "Erving Goffman  ·  1959  ·  Anchor Books",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "Goffman's most famous work introduces DRAMATURGICAL THEORY: social life is essentially a performance. We are all actors on a stage, constantly managing the impressions we give to others.",
                    "He distinguishes FRONT STAGE behaviour (what we present to others — careful, curated, role-appropriate) from BACK STAGE behaviour (how we actually are when nobody is watching).",
                    "Key concepts: impression management (deliberate control of how you appear), face-work (protecting your social dignity and others'), and region behaviour (adapting your performance to the audience and setting).",
                    "Social interaction isn't natural spontaneous expression — it's a skilled, collaborative, moment-by-moment negotiation of meaning and identity.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "PNE's 19 canonical behavioural intentions are performative social acts in Goffman's sense: 'Assert Dominance', 'Deflect with Humour', 'Defend Cause Passionately' are all impression management strategies.",
                    "NPCs in PNE have a public persona (their social layer — assertion, empathy, faction) and an internal cognitive state. The gap between these layers IS Goffman's front-stage/back-stage distinction.",
                    "The choice filter and coherence system ensure NPC dialogue stays in-role — Goffman would call this maintaining a consistent 'performance' to avoid breaking the player's immersion.",
                    "When NPC confrontation_level rises, they 'break character' from their usual front-stage self — this maps to Goffman's concept of frame breaks.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Goffman to justify the separation between NPC's internal BDI state and their expressed dialogue: 'Goffman (1959) characterises social interaction as performance: individuals manage impressions through deliberate presentation of a curated self. PNE encodes this separation structurally — the BDI pipeline determines what the NPC thinks, while the intention layer and LLM determine what they say.'",
                    "This is a powerful point that distinguishes PNE from naive dialogue systems where what the NPC 'feels' IS what they say.",
                    "Also useful in the evaluation chapter: Goffman's dramaturgical lens provides a qualitative framework for assessing whether NPC behaviour feels socially authentic.",
                    "3-5 sentences, positioned at the end of the social psychology section as a capstone.",
                ],
            },
        ],
    },

    # ── Narrative Design ─────────────────────────────────────────────────────
    {
        "category": "Narrative Design",
        "title": "Hamlet on the Holodeck: The Future of Narrative in Cyberspace",
        "author_year": "Janet H. Murray  ·  1997  ·  Free Press",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Book About?",
                "bullets": [
                    "The founding text of interactive narrative studies. Murray asks: what would it look like if the computer became a storytelling medium as rich and culturally significant as the novel or film?",
                    "She identifies three properties that make digital environments unique as narrative spaces: AGENCY (the player makes meaningful choices), IMMERSION (you feel you are inside the story world), and TRANSFORMATION (the story world changes in response to you).",
                    "Murray argues that the challenge of interactive narrative is to give players genuine agency WITHOUT destroying the authorial craft of the story — a tension she calls the 'Holodeck' problem.",
                    "The book is full of examples from early interactive fiction, games, and theory — written before modern open-world games existed, but remarkably prescient.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "PNE is explicitly designed to produce all three of Murray's properties: AGENCY (skill-check choices with real consequence), IMMERSION (LLM-generated natural dialogue grounded in BDI state), TRANSFORMATION (judgement score changes the NPC and routes the story).",
                    "Murray's 'Holodeck problem' is exactly what PNE solves at the dialogue level: how do you give players real choice in NPC conversations without making the NPC incoherent or breaking narrative structure?",
                    "The FSM + judgement-score system is PNE's answer to the agency/authorship tension: players have agency over OUTCOMES, but scenario authors retain control over STRUCTURE.",
                    "Murray's framework gives you a vocabulary to evaluate PNE against — use it explicitly in Chapter 5.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Murray to frame the entire interactive narrative section: 'Murray (1997) identifies agency, immersion, and transformation as the three properties that distinguish digital narrative from static storytelling. These properties frame the design objectives of PNE...'",
                    "Then map each property to a PNE feature: agency → skill-check system; immersion → LLM dialogue generation; transformation → judgement-driven FSM routing.",
                    "Also use Murray's Holodeck problem to motivate WHY PNE is needed: existing dialogue trees (Mass Effect, Witcher) solve authorship but sacrifice agency; procedural systems risk sacrificing narrative coherence.",
                    "Murray should open the narrative design section — it provides the evaluative framework that all other narrative sources then fill in.",
                ],
            },
        ],
    },
    {
        "category": "Narrative Design",
        "title": "Narrative in Virtual Environments — Towards Emergent Narrative",
        "author_year": "Ruth Aylett  ·  1999  ·  AAAI Fall Symposium on Narrative Intelligence",
        "slides": [
            {
                "label": "SLIDE 1 OF 3  ·  THE GIST",
                "heading": "What Is This Paper About?",
                "bullets": [
                    "A landmark 8-page paper that defined the concept of EMERGENT NARRATIVE for the AI and games research community.",
                    "Aylett argues there are two ends of a spectrum: AUTHORED narrative (a fixed story, the player just experiences it) and EMERGENT narrative (characters with goals and beliefs act autonomously — story emerges from their interactions).",
                    "She argues that believable interactive narrative requires moving away from authored scripts toward characters that are genuine autonomous agents with their own internal states.",
                    "The problem: pure emergent narrative is hard to author and hard to guarantee will produce satisfying drama. Aylett acknowledges the tension — interesting drama may require a human shaping hand.",
                ],
            },
            {
                "label": "SLIDE 2 OF 3  ·  PNE APPLICATION",
                "heading": "How Does This Apply to PNE?",
                "bullets": [
                    "PNE sits deliberately in the MIDDLE of Aylett's spectrum: NPCs are autonomous agents (BDI pipeline, emergent responses) BUT within an authored scenario graph (FSM nodes authored by a human).",
                    "This is the key design claim: PNE produces emergent dialogue behaviour (no pre-scripted response) within a structured narrative (authored FSM transitions). Both properties at once.",
                    "Aylett's problem — that pure emergence is hard to control for drama — is exactly why PNE has the judgement score and FSM: the author shapes the story arc, the BDI produces the emergent texture.",
                    "Citing Aylett lets you position PNE as a deliberate resolution to a known problem in the field, not just a feature list.",
                ],
            },
            {
                "label": "SLIDE 3 OF 3  ·  LITERATURE REVIEW",
                "heading": "What to Write in Chapter 2",
                "bullets": [
                    "Use Aylett to define emergent narrative and position PNE within the authored/emergent spectrum: 'Aylett (1999) distinguishes authored narrative, in which story events are predetermined, from emergent narrative, in which story arises from autonomous agent behaviour. PNE occupies a hybrid position...'",
                    "Then explain HOW PNE is hybrid: emergent at the dialogue level (BDI pipeline, LLM generation), authored at the structural level (scenario graph, FSM transitions).",
                    "This is one of the most directly applicable sources to PNE's core design thesis — give it a full paragraph, not just a citation.",
                    "Pair with Murray: Murray defines what good interactive narrative should feel like; Aylett defines the technical problem PNE is solving.",
                ],
            },
        ],
    },
]

# ── Build PPTX ────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs)

    current_category = None
    for paper in PAPERS:
        cat = paper["category"]
        accent = CATEGORY_COLOURS[cat]

        if cat != current_category:
            current_category = cat
            descriptions = {
                "BDI Models":
                    "Three foundational texts on Belief-Desire-Intention agent architectures — the theoretical and practical basis for PNE's cognitive pipeline.",
                "Cognitive Psychology":
                    "Three texts from cognitive psychology that ground PNE's belief formation, cognitive bias modelling, and attribution system in established theory.",
                "Social Psychology":
                    "Three works on persuasion, social identity, and impression management — the theoretical basis for PNE's Language Arts and behavioural intention systems.",
                "Narrative Design":
                    "Two foundational interactive narrative texts: one establishing what good interactive narrative should achieve, one defining the emergent narrative problem PNE is designed to solve.",
            }
            add_section_slide(prs, cat, accent, descriptions.get(cat, ""))

        for s in paper["slides"]:
            add_bullet_slide(
                prs,
                category=cat,
                paper_title=paper["title"],
                author_year=paper["author_year"],
                slide_title=s["heading"],
                bullets=s["bullets"],
                label=s["label"],
                accent_col=accent,
            )

    out_path = os.path.join(
        "d:/Programming/PNE (Github)/cs3ip/Dissertation/Literature Material/Summary",
        "PNE_Literature_Review_Summary.pptx"
    )
    prs.save(out_path)
    print(f"Saved -> {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
